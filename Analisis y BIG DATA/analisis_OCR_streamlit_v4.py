import streamlit as st
import os
import concurrent.futures
import re
import pandas as pd
import math
import warnings
from PIL import Image

import pytesseract
import pdfplumber
import docx
import openpyxl
from odf.opendocument import load
from odf.text import P
import pptx
import csv
import string

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

st.title("Buscador avanzado y visual contextual")

def extrar_texto_pdf(path):
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            with pdfplumber.open(path) as pdf:
                return '\n'.join([page.extract_text() or '' for page in pdf.pages])
    except:
        return ''

def extraer_texto_imagen(path):
    try:
        img = Image.open(path)
        return pytesseract.image_to_string(img)
    except:
        return ''

def extraer_texto_docx(path):
    try:
        doc = docx.Document(path)
        return '\n'.join([para.text for para in doc.paragraphs])
    except:
        return ''

def extraer_texto_odt(path):
    try:
        text = []
        odt = load(path)
        for elem in odt.getElementsByType(P):
            text.append(str(elem))
        return '\n'.join(text)
    except:
        return ''

def extraer_texto_xlsx(path):
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                texts.append(' '.join(str(c) if c is not None else '' for c in row))
        return '\n'.join(texts)
    except:
        return ''

def extraer_texto_txt(path):
    try:
        with open(path, 'r', encoding='utf8', errors='ignore') as f:
            return f.read()
    except:
        return ''

def extraer_texto_pptx(path):
    try:
        prs = pptx.Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return '\n'.join(texts)
    except:
        return ''

def extraer_texto_csv(path):
    try:
        with open(path, 'r', encoding='utf8', errors='ignore') as f:
            reader = csv.reader(f)
            lines = [' '.join(row) for row in reader]
            return '\n'.join(lines)
    except:
        return ''

def extraer_texto_archivo(path):
    ext = path.lower().split('.')[-1]
    if ext == 'pdf':
        return extrar_texto_pdf(path)
    elif ext in ('png', 'jpg', 'jpeg'):
        return extraer_texto_imagen(path)
    elif ext == 'docx':
        return extraer_texto_docx(path)
    elif ext == 'odt':
        return extraer_texto_odt(path)
    elif ext in ('xlsx', 'xls'):
        return extraer_texto_xlsx(path)
    elif ext == 'txt':
        return extraer_texto_txt(path)
    elif ext == 'pptx':
        return extraer_texto_pptx(path)
    elif ext == 'csv':
        return extraer_texto_csv(path)
    else:
        return ''

def buscar_archivos_en_carpeta(carpeta):
    archivos = []
    for root, _, files in os.walk(carpeta):
        for file in files:
            archivos.append(os.path.join(root, file))
    return archivos

def encontrar_fragmento(texto, palabra):
    texto_lower = texto.lower()
    p = palabra.lower()
    idx = texto_lower.find(p)
    if idx == -1:
        return ""
    return texto[max(0, idx-30):idx+70]

def is_ocr_reliable(txt):
    letras = sum([1 for c in txt if c in string.ascii_letters])
    return letras > 20

def resaltar_texto(texto, palabras):
    import html
    texto = html.escape(texto)
    for palabra in palabras:
        palabra_esc = html.escape(palabra)
        patron = re.compile(re.escape(palabra_esc), re.IGNORECASE)
        texto = patron.sub(f'<b>{palabra_esc}</b>', texto)
    return texto

carpeta = st.text_input("Ruta carpeta (se busca en subcarpetas también)").strip().strip('"').strip("'")
entrada_palabras = st.text_input("Palabras a buscar, separadas por coma")

if carpeta and entrada_palabras:
    if not os.path.isdir(carpeta):
        st.error("Ruta no válida o no es carpeta")
    else:
        palabras = [p.strip() for p in entrada_palabras.split(",") if p.strip()]
        archivos = buscar_archivos_en_carpeta(carpeta)
        resultados = []
        progreso = st.progress(0)
        archivo_actual = st.empty()
        resultados_preview = st.empty()

        with st.spinner("Buscando contenido y nombres..."):
            coincidencias_tot = 0
            for i, arch in enumerate(archivos):
                ext = arch.lower().split('.')[-1]
                archivo_actual.text(f"Analizando {i+1} de {len(archivos)}: {os.path.basename(arch)}")
                datos_doc = {'archivo': arch, 'matches': []}

                encontrado_nombre_ruta = False
                for p in palabras:
                    if p.lower() in os.path.basename(arch).lower() or p.lower() in arch.lower():
                        if not encontrado_nombre_ruta:
                            coincidencias_tot += 1
                            encontrado_nombre_ruta = True
                        datos_doc['matches'].append({
                            'tipo': 'nombre/ruta',
                            'palabra': p,
                            'fragmento': os.path.basename(arch)
                        })

                try:
                    contenido = extraer_texto_archivo(arch)
                    if contenido:
                        for p in palabras:
                            ocurrencias = contenido.lower().count(p.lower())
                            if ocurrencias > 0:
                                coincidencias_tot += ocurrencias
                                fragmento = encontrar_fragmento(contenido, p)
                                if ext in ('png', 'jpg', 'jpeg'):
                                    confiable = is_ocr_reliable(contenido)
                                    datos_doc['matches'].append({
                                        'tipo': 'OCR Imagen' + ('' if confiable else ' (NO CONCLUYENTE)'),
                                        'palabra': p,
                                        'fragmento': fragmento if confiable else 'OCR demasiado corto o ruidoso',
                                        'ocr': contenido if confiable else ''
                                    })
                                else:
                                    datos_doc['matches'].append({
                                        'tipo': 'contenido',
                                        'palabra': p,
                                        'fragmento': fragmento
                                    })
                except:
                    pass

                if datos_doc['matches']:
                    resultados.append(datos_doc)

                progreso.progress((i+1)/len(archivos))
                if resultados:
                    tabla = pd.DataFrame([
                        {'Archivo': os.path.basename(r['archivo']),
                         'Tipo': m['tipo'],
                         'Palabra': m['palabra'],
                         'Fragmento': m['fragmento']}
                        for r in resultados for m in r['matches']
                    ])
                    resultados_preview.dataframe(tabla)

        if resultados:
            st.success(f"Búsqueda completada. Total coincidencias (sumando todas): {coincidencias_tot}")
            for res in resultados:
                st.markdown(f"---")
                st.markdown(f"**Archivo:** {res['archivo']}")
                for m in res['matches']:
                    st.write(f"Tipo de coincidencia: {m['tipo']}")
                    st.write(f"Palabra encontrada: {m['palabra']}")
                    st.write(f"Fragmento/contexto:")
                    st.markdown(resaltar_texto(m['fragmento'], [m['palabra']]), unsafe_allow_html=True)
                    ext = res['archivo'].lower().split('.')[-1]
                    if m.get('ocr') and ext in ('png', 'jpg', 'jpeg'):
                        st.info("Texto OCR extraído completo:")
                        st.code(m['ocr'][:800])
                    if ext in ('png','jpg','jpeg'):
                        st.image(Image.open(res['archivo']), caption=res['archivo'])
                with open(res['archivo'], 'rb') as f:
                    st.download_button(label="Descargar", data=f, file_name=os.path.basename(res['archivo']))
        else:
            st.info("No se encontraron coincidencias para esa búsqueda.")
else:
    st.info("Introduce ruta y palabras para buscar.")
